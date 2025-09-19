"""
Servizio per la gestione e elaborazione dei file (documenti, immagini, etc.)
"""

import os
import logging
import hashlib
import mimetypes
from typing import Dict, List, Any, Optional, Tuple, Union
from pathlib import Path
import uuid

# Document processing
import PyPDF2
import fitz  # PyMuPDF
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

# Image processing
from PIL import Image
import pytesseract
import cv2
import numpy as np

logger = logging.getLogger(__name__)

class FileService:
    """Servizio per gestire upload, elaborazione e analisi di file"""
    
    def __init__(self, config: Dict[str, Any]):
        """
        Inizializza il servizio file
        
        Args:
            config: Configurazione del servizio file
        """
        self.config = config
        self.upload_folder = config.get('upload_folder', './data/uploads')
        self.max_file_size = config.get('max_file_size', 50) * 1024 * 1024  # MB to bytes
        self.allowed_extensions = set(config.get('allowed_extensions', [
            'pdf', 'docx', 'txt', 'jpg', 'jpeg', 'png', 'gif', 'bmp', 'xlsx', 'pptx'
        ]))
        
        # Assicura che la directory di upload esista
        os.makedirs(self.upload_folder, exist_ok=True)
        
        logger.info(f"File Service inizializzato - Upload folder: {self.upload_folder}")
    
    def save_uploaded_file(self, file_data: bytes, original_filename: str, 
                          user_id: Optional[str] = None) -> Dict[str, Any]:
        """
        Salva un file caricato e restituisce informazioni
        
        Args:
            file_data: Dati binari del file
            original_filename: Nome originale del file
            user_id: ID dell'utente (opzionale)
            
        Returns:
            Dizionario con informazioni del file salvato
        """
        try:
            # Validazione del file
            validation_result = self.validate_file(file_data, original_filename)
            if not validation_result['valid']:
                return validation_result
            
            # Generazione nome file unico
            file_extension = Path(original_filename).suffix.lower()
            unique_filename = f"{uuid.uuid4().hex}{file_extension}"
            
            # Percorso completo
            file_path = os.path.join(self.upload_folder, unique_filename)
            
            # Salvataggio file
            with open(file_path, 'wb') as f:
                f.write(file_data)
            
            # Calcolo hash del file
            file_hash = self._calculate_file_hash(file_data)
            
            # Informazioni del file
            file_info = {
                'success': True,
                'filename': unique_filename,
                'original_filename': original_filename,
                'file_path': file_path,
                'file_size': len(file_data),
                'file_type': file_extension[1:] if file_extension else 'unknown',
                'mime_type': mimetypes.guess_type(original_filename)[0],
                'file_hash': file_hash,
                'user_id': user_id
            }
            
            logger.info(f"File salvato: {unique_filename} ({len(file_data)} bytes)")
            return file_info
            
        except Exception as e:
            logger.error(f"Errore nel salvataggio del file: {e}")
            return {'success': False, 'error': str(e)}
    
    def validate_file(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """
        Valida un file prima del salvataggio
        
        Args:
            file_data: Dati binari del file
            filename: Nome del file
            
        Returns:
            Risultato della validazione
        """
        try:
            # Controllo dimensione
            if len(file_data) > self.max_file_size:
                return {
                    'valid': False,
                    'error': f'File troppo grande. Massimo {self.max_file_size // (1024*1024)}MB'
                }
            
            # Controllo estensione
            file_extension = Path(filename).suffix.lower()[1:]  # Rimuove il punto
            if file_extension not in self.allowed_extensions:
                return {
                    'valid': False,
                    'error': f'Estensione non supportata. Supportate: {", ".join(self.allowed_extensions)}'
                }
            
            # Controllo contenuto (magic bytes)
            if not self._validate_file_content(file_data, file_extension):
                return {
                    'valid': False,
                    'error': 'Contenuto del file non valido per l\'estensione specificata'
                }
            
            return {'valid': True}
            
        except Exception as e:
            logger.error(f"Errore nella validazione del file: {e}")
            return {'valid': False, 'error': str(e)}
    
    def extract_text_from_document(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """
        Estrae testo da un documento
        
        Args:
            file_path: Percorso del file
            file_type: Tipo del file
            
        Returns:
            Dizionario con testo estratto e metadati
        """
        try:
            if file_type == 'pdf':
                return self._extract_text_from_pdf(file_path)
            elif file_type == 'docx':
                return self._extract_text_from_docx(file_path)
            elif file_type == 'txt':
                return self._extract_text_from_txt(file_path)
            elif file_type == 'xlsx':
                return self._extract_text_from_xlsx(file_path)
            elif file_type == 'pptx':
                return self._extract_text_from_pptx(file_path)
            else:
                return {'success': False, 'error': f'Tipo file non supportato: {file_type}'}
                
        except Exception as e:
            logger.error(f"Errore nell'estrazione testo da {file_path}: {e}")
            return {'success': False, 'error': str(e)}
    
    def extract_text_from_image(self, file_path: str) -> Dict[str, Any]:
        """
        Estrae testo da un'immagine usando OCR
        
        Args:
            file_path: Percorso dell'immagine
            
        Returns:
            Dizionario con testo estratto e metadati
        """
        try:
            # Caricamento immagine
            image = Image.open(file_path)
            
            # Preprocessing dell'immagine per migliorare OCR
            processed_image = self._preprocess_image_for_ocr(image)
            
            # Estrazione testo con Tesseract
            extracted_text = pytesseract.image_to_string(processed_image, lang='ita+eng')
            
            # Informazioni immagine
            image_info = {
                'width': image.width,
                'height': image.height,
                'format': image.format,
                'mode': image.mode
            }
            
            return {
                'success': True,
                'extracted_text': extracted_text.strip(),
                'image_info': image_info,
                'text_length': len(extracted_text.strip())
            }
            
        except Exception as e:
            logger.error(f"Errore nell'estrazione testo da immagine {file_path}: {e}")
            return {'success': False, 'error': str(e)}
    
    def analyze_image(self, file_path: str) -> Dict[str, Any]:
        """
        Analizza un'immagine e estrae informazioni
        
        Args:
            file_path: Percorso dell'immagine
            
        Returns:
            Dizionario con analisi dell'immagine
        """
        try:
            image = Image.open(file_path)
            
            # Informazioni base
            analysis = {
                'success': True,
                'width': image.width,
                'height': image.height,
                'format': image.format,
                'mode': image.mode,
                'size_bytes': os.path.getsize(file_path)
            }
            
            # Analisi colori dominanti
            analysis['dominant_colors'] = self._get_dominant_colors(image)
            
            # Rilevamento oggetti base (se possibile)
            analysis['has_text'] = self._detect_text_in_image(image)
            
            return analysis
            
        except Exception as e:
            logger.error(f"Errore nell'analisi immagine {file_path}: {e}")
            return {'success': False, 'error': str(e)}
    
    def chunk_document_text(self, text: str, chunk_size: int = 1000, 
                           chunk_overlap: int = 200) -> List[str]:
        """
        Divide un testo in chunks per l'elaborazione
        
        Args:
            text: Testo da dividere
            chunk_size: Dimensione di ogni chunk
            chunk_overlap: Sovrapposizione tra chunks
            
        Returns:
            Lista di chunks di testo
        """
        if not text or len(text) <= chunk_size:
            return [text] if text else []
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Cerca un punto di interruzione naturale (spazio, punto, etc.)
            if end < len(text):
                # Cerca l'ultimo spazio o punto nel chunk
                last_space = text.rfind(' ', start, end)
                last_period = text.rfind('.', start, end)
                
                natural_break = max(last_space, last_period)
                if natural_break > start:
                    end = natural_break + 1
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - chunk_overlap
            if start >= len(text):
                break
        
        logger.debug(f"Testo diviso in {len(chunks)} chunks")
        return chunks
    
    def _extract_text_from_pdf(self, file_path: str) -> Dict[str, Any]:
        """Estrae testo da PDF"""
        try:
            # Prova prima con PyMuPDF (più robusto)
            doc = fitz.open(file_path)
            text = ""
            metadata = {
                'page_count': len(doc),
                'title': doc.metadata.get('title', ''),
                'author': doc.metadata.get('author', ''),
                'subject': doc.metadata.get('subject', '')
            }
            
            for page in doc:
                text += page.get_text() + "\n"
            
            doc.close()
            
            return {
                'success': True,
                'extracted_text': text.strip(),
                'metadata': metadata
            }
            
        except Exception as e:
            # Fallback a PyPDF2
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                    
                    metadata = {
                        'page_count': len(pdf_reader.pages),
                        'title': pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                        'author': pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else ''
                    }
                    
                    return {
                        'success': True,
                        'extracted_text': text.strip(),
                        'metadata': metadata
                    }
                    
            except Exception as e2:
                logger.error(f"Errore con entrambi i parser PDF: {e}, {e2}")
                return {'success': False, 'error': str(e2)}
    
    def _extract_text_from_docx(self, file_path: str) -> Dict[str, Any]:
        """Estrae testo da DOCX"""
        doc = DocxDocument(file_path)
        
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Estrae testo dalle tabelle
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        
        # Metadati del documento
        metadata = {
            'title': doc.core_properties.title or '',
            'author': doc.core_properties.author or '',
            'subject': doc.core_properties.subject or '',
            'keywords': doc.core_properties.keywords or ''
        }
        
        return {
            'success': True,
            'extracted_text': text.strip(),
            'metadata': metadata
        }
    
    def _extract_text_from_txt(self, file_path: str) -> Dict[str, Any]:
        """Estrae testo da file TXT"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            text = file.read()
        
        return {
            'success': True,
            'extracted_text': text,
            'metadata': {'encoding': 'utf-8'}
        }
    
    def _extract_text_from_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Estrae testo da Excel"""
        workbook = openpyxl.load_workbook(file_path)
        text = ""
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"\n--- {sheet_name} ---\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text += row_text + "\n"
        
        return {
            'success': True,
            'extracted_text': text.strip(),
            'metadata': {'sheets': workbook.sheetnames}
        }
    
    def _extract_text_from_pptx(self, file_path: str) -> Dict[str, Any]:
        """Estrae testo da PowerPoint"""
        presentation = Presentation(file_path)
        text = ""
        
        for i, slide in enumerate(presentation.slides, 1):
            text += f"\n--- Slide {i} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return {
            'success': True,
            'extracted_text': text.strip(),
            'metadata': {'slide_count': len(presentation.slides)}
        }
    
    def _preprocess_image_for_ocr(self, image: Image.Image) -> Image.Image:
        """Preprocessa un'immagine per migliorare l'OCR"""
        # Converte in array numpy
        img_array = np.array(image)
        
        # Converte in scala di grigi se necessario
        if len(img_array.shape) == 3:
            img_array = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        
        # Applica filtri per migliorare la qualità
        # Denoising
        img_array = cv2.fastNlMeansDenoising(img_array)
        
        # Thresholding per migliorare il contrasto
        _, img_array = cv2.threshold(img_array, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        # Converte di nuovo in PIL Image
        return Image.fromarray(img_array)
    
    def _get_dominant_colors(self, image: Image.Image, num_colors: int = 5) -> List[Tuple[int, int, int]]:
        """Estrae i colori dominanti da un'immagine"""
        # Ridimensiona per velocità
        image = image.resize((150, 150))
        
        # Converte in RGB se necessario
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Ottiene i pixel
        pixels = list(image.getdata())
        
        # Usa k-means clustering per trovare colori dominanti
        from collections import Counter
        
        # Semplificazione: conta i colori più frequenti
        color_counts = Counter(pixels)
        dominant_colors = [color for color, count in color_counts.most_common(num_colors)]
        
        return dominant_colors
    
    def _detect_text_in_image(self, image: Image.Image) -> bool:
        """Rileva se un'immagine contiene testo"""
        try:
            # Estrazione rapida per rilevamento
            text = pytesseract.image_to_string(image, config='--psm 6')
            return len(text.strip()) > 10  # Soglia minima per considerare presenza di testo
        except:
            return False
    
    def _validate_file_content(self, file_data: bytes, file_extension: str) -> bool:
        """Valida il contenuto del file basandosi sui magic bytes"""
        magic_bytes = {
            'pdf': b'%PDF',
            'jpg': b'\xff\xd8\xff',
            'jpeg': b'\xff\xd8\xff',
            'png': b'\x89PNG\r\n\x1a\n',
            'gif': b'GIF8',
            'bmp': b'BM',
            'docx': b'PK\x03\x04',  # ZIP-based formats
            'xlsx': b'PK\x03\x04',
            'pptx': b'PK\x03\x04'
        }
        
        if file_extension in magic_bytes:
            expected_magic = magic_bytes[file_extension]
            return file_data.startswith(expected_magic)
        
        # Per file di testo, controlla che sia decodificabile
        if file_extension == 'txt':
            try:
                file_data.decode('utf-8')
                return True
            except UnicodeDecodeError:
                return False
        
        return True  # Default per estensioni non specificate
    
    def _calculate_file_hash(self, file_data: bytes) -> str:
        """Calcola hash SHA-256 del file"""
        return hashlib.sha256(file_data).hexdigest()
    
    def delete_file(self, file_path: str) -> bool:
        """
        Elimina un file dal filesystem
        
        Args:
            file_path: Percorso del file da eliminare
            
        Returns:
            True se eliminazione riuscita, False altrimenti
        """
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                logger.info(f"File eliminato: {file_path}")
                return True
            else:
                logger.warning(f"File non trovato per eliminazione: {file_path}")
                return False
                
        except Exception as e:
            logger.error(f"Errore nell'eliminazione del file {file_path}: {e}")
            return False
    
    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """
        Ottiene informazioni su un file
        
        Args:
            file_path: Percorso del file
            
        Returns:
            Dizionario con informazioni del file
        """
        try:
            if not os.path.exists(file_path):
                return {'exists': False}
            
            stat = os.stat(file_path)
            
            return {
                'exists': True,
                'size': stat.st_size,
                'created': stat.st_ctime,
                'modified': stat.st_mtime,
                'mime_type': mimetypes.guess_type(file_path)[0]
            }
            
        except Exception as e:
            logger.error(f"Errore nel recupero info file {file_path}: {e}")
            return {'exists': False, 'error': str(e)}
